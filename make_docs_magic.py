from pptx import Presentation
import os
import argparse

def pptx_to_markdown(pptx_file, documentation_dir):
    prs = Presentation(pptx_file)

    # Create Documentation and images directories if they don't exist
    images_dir = os.path.join(documentation_dir, 'images')
    if not os.path.exists(documentation_dir):
        os.makedirs(documentation_dir)
    if not os.path.exists(images_dir):
        os.makedirs(images_dir)

    readme_path = os.path.join(documentation_dir, 'README.md')
    with open(readme_path, 'w') as md:
        for i, slide in enumerate(prs.slides):
            md.write(f'## Slide {i+1}\n\n')
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    md.write(f'{shape.text}\n\n')
                if shape.shape_type == 13:  # Shape type 13 corresponds to a picture
                    image = shape.image
                    image_bytes = image.blob
                    image_filename = f'image_{i}_{shape.name}.png'.replace(' ', '_')  # Replace spaces with underscores
                    image_path = os.path.join(images_dir, image_filename)
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    md.write(f'![Image](./images/{image_filename})\n\n')

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Convert PowerPoint to Markdown.')
    parser.add_argument('pptx_file', help='Path to the PowerPoint file')
    args = parser.parse_args()

    pptx_to_markdown(args.pptx_file, 'Documentation')
